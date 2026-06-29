import io
import csv
import asyncio
import fitz                          # pymupdf  — PDF
import httpx                         # URL fetching
import docx                          # python-docx — .docx
import openpyxl                      # .xlsx
from pptx import Presentation        # .pptx
from bs4 import BeautifulSoup
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from google.oauth2.credentials import Credentials


# ── Google Workspace MIME types → export format ───────────────────────────────
GDRIVE_EXPORT_TYPES = {
    "application/vnd.google-apps.document":     "text/plain",
    "application/vnd.google-apps.spreadsheet":  "text/csv",
    "application/vnd.google-apps.presentation": "text/plain",
}


# ── File upload parser ────────────────────────────────────────────────────────

def parse_document(file_bytes: bytes, filename: str) -> list[str]:
    """Route to the correct parser based on file extension."""
    name = (filename or "").lower()

    if name.endswith(".pdf"):
        return _parse_pdf(file_bytes)
    elif name.endswith(".docx"):
        return _parse_docx(file_bytes)
    elif name.endswith(".xlsx"):
        return _parse_xlsx(file_bytes)
    elif name.endswith(".csv"):
        return _parse_csv(file_bytes)
    elif name.endswith(".pptx"):
        return _parse_pptx(file_bytes)
    elif name.endswith(".html") or name.endswith(".htm"):
        return _parse_html(file_bytes)
    elif name.endswith(".md") or name.endswith(".txt"):
        return _parse_text(file_bytes)
    else:
        try:
            return _parse_pdf(file_bytes)
        except Exception:
            return _parse_text(file_bytes)


def _parse_pdf(file_bytes: bytes) -> list[str]:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    chunks = []
    for page in doc:
        text = page.get_text()
        paragraphs = [p.strip() for p in text.split("\n\n") if len(p.strip()) > 50]
        chunks.extend(paragraphs)
    return chunks


def _parse_docx(file_bytes: bytes) -> list[str]:
    doc = docx.Document(io.BytesIO(file_bytes))
    chunks = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if len(text) > 50:
            chunks.append(text)
    return chunks


def _parse_xlsx(file_bytes: bytes) -> list[str]:
    """
    Each sheet becomes a series of row chunks.
    Rows are serialised as 'Col: value | Col: value' so the LLM
    understands the column context without needing a table renderer.
    """
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
    chunks = []

    for sheet in wb.worksheets:
        rows = list(sheet.iter_rows(values_only=True))
        if not rows:
            continue

        # First row is treated as header
        headers = [str(c).strip() if c is not None else f"Col{i}"
                   for i, c in enumerate(rows[0])]

        for row in rows[1:]:
            # Skip completely empty rows
            if all(c is None for c in row):
                continue
            pairs = [f"{headers[i]}: {str(v)}" for i, v in enumerate(row) if v is not None]
            line = " | ".join(pairs)
            if len(line) > 20:
                chunks.append(f"[Sheet: {sheet.title}] {line}")

    return chunks


def _parse_csv(file_bytes: bytes) -> list[str]:
    """
    Same row-serialisation strategy as XLSX so retrieval works consistently
    across both formats.
    """
    text = file_bytes.decode("utf-8", errors="ignore")
    reader = csv.reader(text.splitlines())
    rows = list(reader)
    if not rows:
        return []

    headers = [h.strip() for h in rows[0]]
    chunks = []

    for row in rows[1:]:
        if not any(row):
            continue
        pairs = [f"{headers[i]}: {v.strip()}" for i, v in enumerate(row)
                 if i < len(headers) and v.strip()]
        line = " | ".join(pairs)
        if len(line) > 20:
            chunks.append(line)

    return chunks


def _parse_pptx(file_bytes: bytes) -> list[str]:
    """
    Extract text from every slide. Each slide becomes one chunk so
    the retrieval context maps cleanly back to individual slides.
    """
    prs = Presentation(io.BytesIO(file_bytes))
    chunks = []

    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            chunks.append(f"[Slide {i}] " + " ".join(texts))

    return chunks


def _parse_html(file_bytes: bytes) -> list[str]:
    """Strip tags and return clean paragraphs — same logic as the URL scraper."""
    soup = BeautifulSoup(file_bytes, "lxml")
    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()
    content = soup.find("article") or soup.find("main") or soup.body or soup
    raw = content.get_text(separator="\n")
    return [p.strip() for p in raw.split("\n") if len(p.strip()) > 50]


def _parse_text(file_bytes: bytes) -> list[str]:
    text = file_bytes.decode("utf-8", errors="ignore")
    paragraphs = [p.strip() for p in text.split("\n\n") if len(p.strip()) > 50]
    return paragraphs


# ── URL scraper ───────────────────────────────────────────────────────────────

async def parse_url(url: str) -> list[str]:
    """Fetch a public URL and return clean text paragraphs."""
    async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client:
        response = await client.get(url, headers={"User-Agent": "RAGForge/1.0"})
        response.raise_for_status()

    soup = BeautifulSoup(response.text, "lxml")

    for tag in soup(["script", "style", "nav", "footer", "header", "aside"]):
        tag.decompose()

    content = soup.find("article") or soup.find("main") or soup.body or soup
    raw = content.get_text(separator="\n")
    paragraphs = [p.strip() for p in raw.split("\n") if len(p.strip()) > 50]
    return paragraphs


# ── Google Drive fetcher ──────────────────────────────────────────────────────

async def parse_gdrive(file_id: str, access_token: str) -> list[str]:
    return await asyncio.to_thread(_parse_gdrive_sync, file_id, access_token)


def _parse_gdrive_sync(file_id: str, access_token: str) -> list[str]:
    creds = Credentials(token=access_token)
    service = build("drive", "v3", credentials=creds, cache_discovery=False)

    meta = service.files().get(fileId=file_id, fields="name,mimeType").execute()
    mime = meta.get("mimeType", "")
    name = meta.get("name", "")

    buffer = io.BytesIO()

    if mime in GDRIVE_EXPORT_TYPES:
        export_mime = GDRIVE_EXPORT_TYPES[mime]
        request = service.files().export_media(fileId=file_id, mimeType=export_mime)
        name = name + ".txt"
    else:
        request = service.files().get_media(fileId=file_id)

    downloader = MediaIoBaseDownload(buffer, request)
    done = False
    while not done:
        _, done = downloader.next_chunk()

    return parse_document(buffer.getvalue(), name)